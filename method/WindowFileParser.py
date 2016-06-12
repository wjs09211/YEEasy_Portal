# -*- coding: utf-8 -*-

from pptx import Presentation
import urllib
import chardet
import os
from docx import Document
import PyPDF2
import subprocess
from Info.TextColors import Color


def get_unicode(key):
    encode = chardet.detect(key)['encoding']
    if encode == 'windows-1252':
        encode = 'Big5'
    return unicode(key, encode)


def pdf_parser(file_path, key):
    try:
        key = get_unicode(key)
        pdf_file = open(file_path, 'rb')
        pdf_reader = PyPDF2.PdfFileReader(pdf_file)

        lines = []
        for i in range(pdf_reader.numPages):
            page = pdf_reader.getPage(i)
            if key in page.extractText():
                lines.append(i + 1)

        if len(lines) != 0:
            print Color.BLUE + "filename:" + Color.ENDC, os.path.basename(file_path)
            print Color.BLUE + "in page:" + Color.ENDC, lines
            print
            return True
        else:
            return False
    except:
        print u'encode pdf error'


def doc_parser(file_path, key):
    try:
        key = get_unicode(key)
        p = subprocess.Popen("antiword \"" + file_path + "\" | grep -in " + key + " | awk -F: \'{print $1}\'", stdout=subprocess.PIPE, shell=True)
        (output, err) = p.communicate()
        lines = output.split('\n')[:-1]

        if len(lines) != 0:
            print Color.BLUE + "filename:" + Color.ENDC, os.path.basename(file_path)
            print Color.BLUE + "in Line:" + Color.ENDC, lines
            print
            return True
        else:
            return False
    except:
        print u'encode doc error'


def docx_parser(file_path, key):
    try:
        key = get_unicode(key)
        document = Document(file_path)
        lines = []
        for i, para in enumerate(document.paragraphs):
            if key in para.text:
                lines.append(i + 1)

        if len(lines) != 0:
            print Color.BLUE + "filename:" + Color.ENDC, os.path.basename(file_path)
            print Color.BLUE + "in Line:" + Color.ENDC, lines
            print
            return True
        else:
            return False
    except:
        print u'encode docx error'


def pptx_parser(file_path, key):
    try:
        key = get_unicode(key)
        prs = Presentation(file_path)
        lines = []

        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.find(key) != -1:
                                if i not in lines:
                                    lines.append(i)
        if len(lines) != 0:
            print Color.BLUE + "filename:" + Color.ENDC, os.path.basename(file_path)
            print Color.BLUE + "in slide:" + Color.ENDC, lines
            print
            return True
        else:
            return False
    except:
        print u'encode pptx error'


def txt_parser(file_path, key):
    try:
        key = get_unicode(key)
        with open(file_path, 'rb') as f:
            content = f.read()
        content = content.decode(chardet.detect(content)['encoding'], 'ignore')

        lines = []
        for i, line in enumerate(content.split('\n')):
            if key in line:
                lines.append(i + 1)
        if len(lines) != 0:
            print Color.BLUE + "filename:" + Color.ENDC, os.path.basename(file_path)
            print Color.BLUE + "in Line:" + Color.ENDC, lines
            print
            return True
        else:
            return False
    except:
        print u"encode txt error"


if __name__ == "__main__":
    # txt_parser(u"material/unix+final+project+proposal+comments.txt", u"分數")
    # txt_parser(u"material/unix課堂中proposal討論分數_20160504.txt", u"分數")
    s = "分數"
    print unicode(s, 'utf-8')

