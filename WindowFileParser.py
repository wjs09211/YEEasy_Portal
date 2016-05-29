# -*- coding: utf-8 -*-

from pptx import Presentation
import urllib
import chardet
import os
from docx import Document


def docx_parser(file_path, key):
    encode = chardet.detect(key)['encoding']
    if encode == 'windows-1252':
        encode = 'Big5'
    key = unicode(key, encode)
    document = Document(file_path)
    lines = []
    for i, para in enumerate(document.paragraphs):
        if key in para.text:
            lines.append(i + 1)
    print lines

def ppt_parser(file_path, key):
    try:
        encode = chardet.detect(key)['encoding']
        if encode == 'windows-1252':
            encode = 'Big5'
        key = unicode(key, encode)

        prs = Presentation(file_path)
        lines = []

        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.find(key) != -1:
                                if i not in lines:
                                    lines.append(i)
        if len(lines) != 0:
            print "file_name:", os.path.basename(file_path)
            print "in slide:", lines
            print
            return True
        else:
            return False
    except:
        print u'encode something error'


def txt_parser(file_path, key):
    try:
        encode = chardet.detect(key)['encoding']
        if encode == 'windows-1252':
            encode = 'Big5'
        key = unicode(key, encode)

        with open(file_path, 'rb') as f:
            content = f.read()
        content = content.decode(chardet.detect(content)['encoding'], 'ignore')

        lines = []
        for i, line in enumerate(content.split('\n')):
            if key in line:
                lines.append(i + 1)
        if len(lines) != 0:
            print "filename:", os.path.basename(file_path)
            print "key world:", key
            print "in Line:", lines
            print
            return True
        else:
            return False
    except:
        print u"encode something error"


if __name__ == "__main__":
    # txt_parser(u"material/unix+final+project+proposal+comments.txt", u"分數")
    # txt_parser(u"material/unix課堂中proposal討論分數_20160504.txt", u"分數")
    s = "分數"
    print unicode(s, 'utf-8')

